from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# --- Helpers --------------------------------------------------------------
def add_slide(prs, title, body_lines=None, notes=None, subtitle=None):
    """
    body_lines: list of strings. Indent with "  - " for sub-bullets.
    """
    # Title+Content layout
    layout = prs.slide_layouts[1]  # 0=Title, 1=Title+Content
    slide = prs.slides.add_slide(layout)

    # Title
    slide.shapes.title.text = title

    # Optional subtitle (we'll put it as the first paragraph in body if provided)
    if subtitle and (body_lines is None):
        body_lines = [subtitle]
    elif subtitle:
        body_lines = [subtitle] + (body_lines or [])

    # Body
    if body_lines:
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_lines):
            level = 0
            text = line
            # simple indent convention: two spaces + hyphen
            if line.lstrip().startswith("- "):
                text = line.lstrip()[2:]
            if line.startswith("  - "):
                level = 1
                text = line[4:]
            if line.startswith("    - "):
                level = 2
                text = line[6:]

            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = text
            p.level = level
            p.font.size = Pt(20)

    # Notes
    if notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes

    return slide

def add_table_slide(prs, title, notes, header, rows):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    # Insert a small intro above the table
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    intro = body.paragraphs[0]
    intro.text = ""
    intro.space_after = Pt(6)

    # Add table
    rows_n = len(rows) + 1
    cols_n = len(header)
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(8.4)
    height = Inches(0.8 + 0.3 * rows_n)
    table = slide.shapes.add_table(rows_n, cols_n, left, top, width, height).table

    # Header
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True

    # Rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)

    if notes:
        notes_tf = slide.notes_slide.notes_text_frame
        notes_tf.clear()
        notes_tf.text = notes

    return slide

# --- Build deck -----------------------------------------------------------
prs = Presentation()

# Simple, readable font sizing
for s in prs.slide_layouts:
    pass  # using defaults; we’ll set sizes on bullets

# Slide 1 — Title
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
slide.shapes.title.text = "GitHub Copilot Enablement at Scale"
subtitle = slide.placeholders[1]
subtitle.text = "From Pilot to Policy\nJohn [Last Name] – AI Enablement & Automation Engineer"
slide.notes_slide.notes_text_frame.text = (
    "Introduce yourself and frame the story as a structured, data-driven rollout."
)

# Slide 2 — The Challenge
add_slide(
    prs,
    "The Challenge",
    body_lines=[
        "Pain Points:",
        "- Developers overwhelmed by boilerplate",
        "- Uneven prompt quality",
        "- Leadership anxious about AI-generated code",
        "",
        "Goals:",
        "- Structure prompt training",
        "- Define security guardrails",
        "- Quantify ROI",
    ],
    notes="Emphasize this was about trust, not just productivity."
)

# Slide 3 — The Pilot
add_slide(
    prs,
    "Pilot Overview",
    body_lines=[
        "- 10 engineers across automation and DevOps",
        "- 4-week timeline using real Jira tickets",
        "- Deliverables: ETL jobs, Ansible roles, Jenkins tasks",
    ],
    notes="Highlight that it was production work, not a demo."
)

# Slide 4 — Framework
add_slide(
    prs,
    "Framework for Safe Acceleration",
    body_lines=[
        "Prompt → Review → Measure → Improve → Scale",
        "",
        "Four Pillars:",
        "- Prompt Playbook",
        "- Review Checklist",
        "- Metrics Tracker",
        "- Guardrails",
    ],
    notes="Explain the loop created continuous improvement and measurable progress."
)

# Slide 5 — Prompt Playbook Example
add_slide(
    prs,
    "Prompt Playbook Example",
    body_lines=[
        "Naïve: “Write a Python script to parse logs.”",
        "Structured: “Write portable parse_logs.py that accepts --file, counts ERROR lines, outputs a JSON summary.”",
    ],
    notes="Show how structured prompting directly improves accuracy and consistency."
)

# Slide 6 — Governance & Security
add_slide(
    prs,
    "Governance & Security",
    body_lines=[
        "- No proprietary data in prompts",
        "- Mandatory PR review for Copilot output",
        "- SonarQube + pre-commit scanning",
        "- Tag AI-assisted commits",
    ],
    notes="Present this as enabling responsible creativity, not limiting it."
)

# Slide 7 — Quantitative Results (table)
add_table_slide(
    prs,
    "Quantitative Results",
    notes="Emphasize focus on real business metrics—speed, quality, morale.",
    header=["Metric", "Baseline", "With Copilot", "Δ"],
    rows=[
        ["Routine script time", "60 min", "25 min", "↓ 58%"],
        ["PR turnaround", "2.3 days", "1.5 days", "↓ 35%"],
        ["Developer satisfaction", "—", "8.7 / 10", "↑"],
    ]
)

# Slide 8 — Developer Voices
add_slide(
    prs,
    "Developer Voices",
    body_lines=[
        "“It feels like pair-programming with documentation built in.”",
        "“I stopped tabbing to StackOverflow 50 times a day.”",
        "“I’m learning faster than ever.”",
    ],
    notes="Authentic voices from engineers prove real adoption."
)

# Slide 9 — Lessons Learned
add_slide(
    prs,
    "Lessons Learned",
    body_lines=[
        "- Train before rollout — skill before scale",
        "- Treat prompts as code — version and review them",
        "- Govern lightly — empower, don’t police",
        "- Automate ROI tracking — data drives trust",
    ],
    notes="Stress that enablement is cultural as well as technical."
)

# Slide 10 — Scaling to Enterprise
add_slide(
    prs,
    "Scaling to Enterprise",
    body_lines=[
        "- 25-user pilot",
        "- Apply playbooks + metrics templates",
        "- Extend guardrails to M365 Copilot",
        "- Integrate ROI tracking into CI/CD",
    ],
    notes="Offer this as Duke Energy’s path to safe enterprise adoption."
)

# Slide 11 — The Payoff
add_slide(
    prs,
    "The Payoff",
    body_lines=[
        "Structured onboarding + measurable outcomes = sustainable adoption.",
        "“Responsible AI doesn’t slow teams—it lets them accelerate safely.”",
    ],
    notes="End with confidence and an open invitation for discussion."
)

prs.save("copilot_rollout_deck.pptx")
print("Wrote: copilot_rollout_deck.pptx")
